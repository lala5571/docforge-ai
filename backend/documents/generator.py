from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document as DocxDocument
from docx.shared import Pt as DocPt
import os
from django.conf import settings

def create_pptx(content, document_id):
    prs = Presentation()
    
    # Slide dimensions
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for slide_data in content['slides']:
        # Blank slide layout
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Background — dark blue
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(6, 8, 16)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(12), Inches(1.2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data['title']
        title_para.runs[0].font.size = Pt(32)
        title_para.runs[0].font.bold = True
        title_para.runs[0].font.color.rgb = RGBColor(108, 99, 255)

        # Content points
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(12), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, point in enumerate(slide_data['content']):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            para.text = f"• {point}"
            para.runs[0].font.size = Pt(20)
            para.runs[0].font.color.rgb = RGBColor(240, 242, 255)
            para.space_after = Pt(12)

    # Save file
    os.makedirs(f"{settings.MEDIA_ROOT}/clean", exist_ok=True)
    filepath = f"{settings.MEDIA_ROOT}/clean/{document_id}.pptx"
    prs.save(filepath)
    return filepath


def create_docx(content, document_id):
    doc = DocxDocument()

    # Title
    title = doc.add_heading(content['title'], 0)
    title.runs[0].font.color.rgb = RGBColor(108, 99, 255)

    # Sections
    for section in content['sections']:
        doc.add_heading(section['heading'], level=1)
        doc.add_paragraph(section['content'])
        doc.add_paragraph()

    # Save file
    os.makedirs(f"{settings.MEDIA_ROOT}/clean", exist_ok=True)
    filepath = f"{settings.MEDIA_ROOT}/clean/{document_id}.docx"
    doc.save(filepath)
    return filepath